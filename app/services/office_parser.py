from docx import Document as DocxDocument
from pptx import Presentation


def extract_docx_text(path: str) -> str:
    document = DocxDocument(path)

    sections = [p.text for p in document.paragraphs if p.text.strip()]

    for table in document.tables:
        rows = ["\t".join(cell.text for cell in row.cells) for row in table.rows]
        if rows:
            sections.append("\n".join(rows))

    return "\n\n".join(sections)


def extract_pptx_text(path: str) -> str:
    presentation = Presentation(path)

    sections = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_lines.append(shape.text_frame.text)
            elif shape.has_table:
                rows = ["\t".join(cell.text for cell in row.cells) for row in shape.table.rows]
                slide_lines.extend(rows)
        if slide_lines:
            sections.append(f"Слайд {slide_index}:\n" + "\n".join(slide_lines))

    return "\n\n".join(sections)
