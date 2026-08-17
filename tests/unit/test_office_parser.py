from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches

from app.services.office_parser import extract_docx_text, extract_pptx_text


def _make_docx(tmp_path, paragraphs: list[str], table_rows: list[list[str]] | None = None) -> str:
    document = DocxDocument()
    for text in paragraphs:
        document.add_paragraph(text)
    if table_rows:
        table = document.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        for row_idx, row in enumerate(table_rows):
            for col_idx, cell_text in enumerate(row):
                table.cell(row_idx, col_idx).text = cell_text
    path = tmp_path / "test.docx"
    document.save(path)
    return str(path)


def _make_pptx(tmp_path, slide_texts: list[str]) -> str:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    for text in slide_texts:
        slide = presentation.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        textbox.text_frame.text = text
    path = tmp_path / "test.pptx"
    presentation.save(path)
    return str(path)


def test_extract_docx_text_returns_paragraphs(tmp_path):
    path = _make_docx(tmp_path, ["Первый абзац про модуль P2.5", "Второй абзац"])

    text = extract_docx_text(path)

    assert "Первый абзац про модуль P2.5" in text
    assert "Второй абзац" in text


def test_extract_docx_text_includes_tables(tmp_path):
    path = _make_docx(tmp_path, ["Заголовок"], table_rows=[["Модуль", "Мощность"], ["P2.5", "45Вт"]])

    text = extract_docx_text(path)

    assert "Модуль" in text
    assert "P2.5" in text
    assert "45Вт" in text


def test_extract_docx_text_skips_empty_paragraphs(tmp_path):
    path = _make_docx(tmp_path, ["Текст", "", "  "])

    text = extract_docx_text(path)

    assert text.strip() == "Текст"


def test_extract_pptx_text_returns_slide_text_with_slide_numbers(tmp_path):
    path = _make_pptx(tmp_path, ["Первый слайд про LED-экраны", "Второй слайд"])

    text = extract_pptx_text(path)

    assert "Слайд 1" in text
    assert "Первый слайд про LED-экраны" in text
    assert "Слайд 2" in text
    assert "Второй слайд" in text


def test_extract_pptx_text_empty_presentation_returns_empty_string(tmp_path):
    path = tmp_path / "empty.pptx"
    Presentation().save(path)

    text = extract_pptx_text(str(path))

    assert text == ""
